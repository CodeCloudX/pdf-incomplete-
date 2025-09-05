# tools/pdf_to_ppt_tool.py
import os
import time
import logging
from flask import current_app
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt
from utils.file_utils import validate_file_size
from utils.file_manager import get_session_folder
from utils.file_naming_utils import generate_file_names

logger = logging.getLogger(__name__)

def pdf_to_ppt(file_path, pages=None, slide_width=10.0, slide_height=7.5):
    """
    Convert PDF pages to PowerPoint slides.
    pages: list of 1-based page numbers to convert (optional)
    slide_width, slide_height: slide dimensions in inches
    """
    try:
        if not validate_file_size(file_path):
            return {
                'status': 'error',
                'output_files': [],
                'message': f"File {os.path.basename(file_path)} exceeds size limit"
            }

        # Get original filename for display purposes
        original_filename = os.path.basename(file_path)
        
        # Generate secure file names for output
        file_names = generate_file_names(original_filename, toolname='ppt', ext='pptx')
        display_name = file_names['display_name']
        stored_name = file_names['stored_name']

        # Get session-specific folders
        processed_folder = get_session_folder('processed')
        upload_folder = get_session_folder('uploads')
        out_path = os.path.join(processed_folder, stored_name)
        
        # Ensure directories exist
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        os.makedirs(upload_folder, exist_ok=True)

        # Convert PDF pages to images
        if pages:
            first_page = min(pages)
            last_page = max(pages)
            images = convert_from_path(file_path, first_page=first_page, last_page=last_page)
        else:
            images = convert_from_path(file_path)

        if not images:
            return {
                'status': 'error',
                'output_files': [],
                'message': "No pages found for conversion"
            }

        prs = Presentation()
        prs.slide_width = Inches(slide_width)
        prs.slide_height = Inches(slide_height)

        temp_images = []
        
        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            img_path = os.path.join(upload_folder, f"temp_slide_{time.time()}_{stored_name}.png")
            img.save(img_path, "PNG")
            temp_images.append(img_path)
            
            left = top = 0
            slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)

        # Save PowerPoint
        prs.save(out_path)

        # Cleanup temporary images
        for img_path in temp_images:
            if os.path.exists(img_path):
                os.remove(img_path)

        logger.info(f"Converted {file_path} to PowerPoint: {out_path}")
        return {
            'status': 'success',
            'output_files': [{
                'display_name': f"ppt_{display_name}.pptx",
                'stored_name': stored_name,
                'output_path': out_path
            }],
            'message': f"Converted {len(images)} pages to PowerPoint slides"
        }

    except Exception as e:
        logger.error(f"PDF to PPT conversion failed for {file_path}: {str(e)}")
        
        # Cleanup any temporary files on error
        upload_folder = get_session_folder('uploads')
        temp_files = [f for f in os.listdir(upload_folder) if f.startswith(f"temp_slide_") and stored_name in f]
        for temp_file in temp_files:
            temp_path = os.path.join(upload_folder, temp_file)
            if os.path.exists(temp_path):
                os.remove(temp_path)
                
        return {
            'status': 'error',
            'output_files': [],
            'message': f"Conversion failed: {str(e)}"
        }
